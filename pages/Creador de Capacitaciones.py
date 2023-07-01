import streamlit as st
import base64
import glob
import os
import random
import re
import string
from urllib.parse import urlparse

import openai
from icrawler import ImageDownloader
from icrawler.builtin import GoogleImageCrawler
from pptx import Presentation


# Estableciendo el logo de ACHS
st.image("img/logo_mundo.png", width=300)

# Título y descripción
st.subheader("Asistente Creador de Presentaciones PPT")
st.markdown("""
<p style='font-size: 15px;'>
Este es un asistente virtual diseñado para ayudarte a crear presentaciones en PowerPoint (PPT). 
Simplemente ingresa tu contenido y la cantidad de diapositivas que deseas, y yo me encargaré de generar una presentación visualmente atractiva.

Ej: Crea una presentación sobre los riesgos laborales
</p>
""", unsafe_allow_html=True)

# Define las variables de entrada aquí

topic = st.text_input("Ingresa tu solicitud")  
slide_length = st.slider("Cantidad de páginas sugeridas", min_value=1, max_value=15, value=5)

# Definiendo la API
api_key = os.getenv("OPENAI_API_KEY")

unique_image_name = None

def main():
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                                range(16))

    def refresh_unique_image_name():
        global unique_image_name
        unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                    for _ in range(16))
        return

    class PrefixNameDownloader(ImageDownloader):

        def get_filename(self, task, default_ext):
            url_path = urlparse(task['file_url'])[2]
            if '.' in url_path:
                extension = url_path.split('.')[-1]
                if extension.lower() not in [
                    'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'ppm', 'pgm'
                ]:
                    extension = default_ext
            else:
                extension = default_ext
            st.write(unique_image_name)
            filename = base64.b64encode(url_path.encode()).decode()
            return "p_" + unique_image_name + '{}.{}'.format(filename, extension)

    # Mover la función generate_ppt aquí
    def generate_ppt(topic, slide_length, api_key):
        root = Presentation("theme0.pptx")

        openai.api_key = api_key

        message = f"""Crear un esquema para una presentación de diapositivas sobre el tema de {topic} cual es {slide_length}
        diapositivas largas. Asegúrate de que sea {slide_length} de largo.

        Puede utilizar los siguientes tipos de diapositivas:
        Diapositiva de título - (Title, Subtitle)
        Diapositiva de contenido: (Title, Content)
        Diapositiva de imagen: (Title, Content, Image)
        Diapositiva de agradecimiento -(Title)

        Coloque esta etiqueta antes de la diapositiva de título: [L_TS]
        Coloque esta etiqueta antes de la diapositiva de contenido: [L_CS]
        Coloque esta etiqueta antes de la diapositiva de imagen: [L_IS]
        Coloque esta etiqueta antes de la diapositiva de agradecimiento: [L_THS]
        
        Pon esta etiqueta antes del Título: [TITLE]
        Ponga esta etiqueta después del Título: [/TITLE]
        Pon esta etiqueta antes del Subtítulo: [SUBTITLE]
        Pon esta etiqueta después del Subtítulo: [/SUBTITLE]
        Pon esta etiqueta antes del Contenido: [CONTENT]
        Pon esta etiqueta después del Contenido: [/CONTENT]
        Pon esta etiqueta antes de la Imagen: [IMAGE]
        Pon esta etiqueta después de la Imagen: [/IMAGE]

        Pon "[SLIDEBREAK]" después de cada diapositiva

        Por ejemplo:
        [L_TS]
        [TITLE]Entre nosotros[/TITLE]

        [SLIDEBREAK]

        [L_CS]
        [TITLE]¿Qué hay entre nosotros?[/TITLE]
        [CONTENT]
        1. Among Us es un popular juego multijugador en línea desarrollado y publicado por InnerSloth.
        2. El juego se desarrolla en un entorno espacial donde los jugadores asumen el papel de compañeros de tripulación e impostores. 
        3. El objetivo de los Crewmates es completar tareas e identificar a los impostores entre ellos, mientras que el objetivo de los impostores es sabotear la nave espacial y eliminar a los Crewmates sin ser atrapados.
        [/CONTENT]

        [SLIDEBREAK]


        Elaborar el contenido, proporcionar la mayor cantidad de información posible.
        RECUERDA COLOCAR un [/CONTENT] al final del Contenido.
        No incluya caracteres especiales (?, !, ., :, ) en el Título.
        No incluya ninguna información adicional en su respuesta y respete el formato."""

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": message}
            ]
        )

        # """ Ref for slide types:
        # 0 -> title and subtitle
        # 1 -> title and content
        # 2 -> section header
        # 3 -> two content
        # 4 -> Comparison
        # 5 -> Title only
        # 6 -> Blank
        # 7 -> Content with caption
        # 8 -> Pic with caption
        # """

        def delete_all_slides():
            for i in range(len(root.slides) - 1, -1, -1):
                r_id = root.slides._sldIdLst[i].rId
                root.part.drop_rel(r_id)
                del root.slides._sldIdLst[i]

        def create_title_slide(title, subtitle):
            layout = root.slide_layouts[0]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = subtitle

        def create_section_header_slide(title):
            layout = root.slide_layouts[2]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title

        def create_title_and_content_slide(title, content):
            layout = root.slide_layouts[1]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        def create_title_and_content_and_image_slide(title, content, image_query):
            layout = root.slide_layouts[8]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[2].text = content
            refresh_unique_image_name()
            google_crawler = GoogleImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': os.getcwd()})
            google_crawler.crawl(keyword=image_query, max_num=1)
            file_name = glob.glob(f"p_{unique_image_name}*")
            if file_name:
                img_path = file_name[0]
                slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                        slide.placeholders[1].width, slide.placeholders[1].height)
            else:
                st.write("Error: No se encontró la imagen descargada")


        def find_text_in_between_tags(text, start_tag, end_tag):
            start_pos = text.find(start_tag)
            end_pos = text.find(end_tag)
            result = []
            while start_pos > -1 and end_pos > -1:
                text_between_tags = text[start_pos + len(start_tag):end_pos]
                result.append(text_between_tags)
                start_pos = text.find(start_tag, end_pos + len(end_tag))
                end_pos = text.find(end_tag, start_pos)
            res1 = "".join(result)
            res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
            if len(result) > 0:
                return res2
            else:
                return ""

        def search_for_slide_type(text):
                tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
                found_text = next((s for s in tags if s in text), None)
                return found_text

        def parse_response(reply):
            list_of_slides = reply.split("[SLIDEBREAK]")
            for slide in list_of_slides:
                slide_type = search_for_slide_type(slide)
                if slide_type == "[L_TS]":
                    create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                    find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                elif slide_type == "[L_CS]":
                    create_title_and_content_slide(
                        "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                        "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                        "[/CONTENT]")))
                elif slide_type == "[L_IS]":
                    create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                            "[/TITLE]")),
                                                            "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                            "[/CONTENT]")),
                                                            "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                            "[/IMAGE]")))
                elif slide_type == "[L_THS]":
                    create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        def find_title():
            title = root.slides[0].shapes.title.text
            if title:
                return title
            else:
                return "Untitled"

        delete_all_slides()

        print(response)

        parse_response(response['choices'][0]['message']['content'])

        root.save(f"{find_title()}.pptx")

        print("done")

        return rf"Hecho {find_title()} está listo, Puedes encontrarlo en {os.getcwd()}\{find_title()}.pptx"

    if st.button("Generar"):
        st.write(generate_ppt(topic, slide_length, api_key))

main()


# Estableciendo la franja superior
st.image("img/franja_inferior.png")
